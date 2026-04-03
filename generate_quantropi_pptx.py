#!/usr/bin/env python3
"""
Generate a Quantropi-branded PowerPoint template for Randy.
Brand Guidelines:
  - Theme: Light
  - Primary: #C7AADC, Accent: #7BA4DD, Background: #FFFFFF
  - Text Primary: #101743
  - Font: Source Sans Pro (headings + body)
  - Tone: Professional, Energy: Medium
  - Audience: Enterprise clients seeking quantum security solutions
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──────────────────────────────────────────────
PRIMARY    = RGBColor(0xC7, 0xAA, 0xDC)   # Lavender purple
ACCENT     = RGBColor(0x7B, 0xA4, 0xDD)   # Soft blue
BG_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK   = RGBColor(0x10, 0x17, 0x43)   # Dark navy
LIGHT_GRAY = RGBColor(0xF5, 0xF3, 0xF8)   # Very light purple tint
MID_GRAY   = RGBColor(0x6B, 0x70, 0x8D)   # Subtle text gray
PRIMARY_LIGHT = RGBColor(0xE8, 0xDD, 0xF0) # Very light primary

# ── Font ──────────────────────────────────────────────────────
FONT_PRIMARY = "Source Sans Pro"

# ── Slide dimensions (16:9) ───────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ── Helper Functions ──────────────────────────────────────────

def add_bg(slide, color=BG_WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Adjust rounding
    shape.adjustments[0] = 0.05
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=20,
                color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_PRIMARY, anchor=MSO_ANCHOR.TOP):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Set vertical alignment
    try:
        tf.paragraphs[0].alignment = alignment
    except:
        pass

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple styled lines.
    lines: list of dicts with keys: text, font_size, color, bold, spacing_after
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line.get("text", "")
        p.font.size = Pt(line.get("font_size", 20))
        p.font.color.rgb = line.get("color", TEXT_DARK)
        p.font.bold = line.get("bold", False)
        p.font.name = line.get("font_name", FONT_PRIMARY)
        p.alignment = alignment
        if "spacing_after" in line:
            p.space_after = Pt(line["spacing_after"])
        if "spacing_before" in line:
            p.space_before = Pt(line["spacing_before"])

    return txBox


def add_decorative_bar(slide, left, top, width, height, color):
    """Add a thin decorative accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color, opacity=1.0):
    """Add a circle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_gradient_bar(slide, left, top, width, height):
    """Add a gradient bar from PRIMARY to ACCENT."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = PRIMARY
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = ACCENT
    fill.gradient_stops[1].position = 1.0
    shape.line.fill.background()
    return shape


def add_footer_bar(slide):
    """Add standard footer with branding."""
    # Thin gradient line
    add_gradient_bar(slide, Inches(0.8), Inches(6.8), Inches(11.7), Pt(2))
    # Footer text
    add_textbox(slide, Inches(0.8), Inches(6.9), Inches(5), Inches(0.4),
                "Quantropi  |  Quantum Security Solutions", font_size=10,
                color=MID_GRAY, bold=False)
    add_textbox(slide, Inches(9.5), Inches(6.9), Inches(3), Inches(0.4),
                "CONFIDENTIAL", font_size=10,
                color=MID_GRAY, bold=False, alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ══════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_bg(slide1, BG_WHITE)

# Left decorative panel
add_rect(slide1, Inches(0), Inches(0), Inches(0.4), SLIDE_H, PRIMARY)

# Decorative circles (subtle brand elements)
add_circle(slide1, Inches(9.5), Inches(0.5), Inches(2.5), PRIMARY_LIGHT)
add_circle(slide1, Inches(10.8), Inches(1.2), Inches(1.8), RGBColor(0xD5, 0xE3, 0xF3))

# Main title area
add_textbox(slide1, Inches(1.2), Inches(1.5), Inches(8), Inches(0.5),
            "QUANTROPI", font_size=14, color=PRIMARY, bold=True)

add_decorative_bar(slide1, Inches(1.2), Inches(2.1), Inches(1.5), Pt(4), PRIMARY)

add_textbox(slide1, Inches(1.2), Inches(2.5), Inches(8), Inches(1.5),
            "Presentation Title", font_size=40, color=TEXT_DARK, bold=True)

add_textbox(slide1, Inches(1.2), Inches(4.2), Inches(8), Inches(0.8),
            "Subtitle or event description goes here", font_size=20,
            color=MID_GRAY, bold=False)

# Presenter info block
add_rounded_rect(slide1, Inches(1.2), Inches(5.3), Inches(4.5), Inches(1.2), LIGHT_GRAY)
add_multiline_textbox(slide1, Inches(1.5), Inches(5.45), Inches(4), Inches(1),
    [
        {"text": "Presenter Name", "font_size": 16, "bold": True, "color": TEXT_DARK, "spacing_after": 4},
        {"text": "Title  |  Quantropi Inc.", "font_size": 12, "color": MID_GRAY, "spacing_after": 2},
        {"text": "Date", "font_size": 12, "color": MID_GRAY},
    ])

# Bottom gradient bar
add_gradient_bar(slide1, Inches(0), Inches(7.3), SLIDE_W, Pt(6))


# ══════════════════════════════════════════════════════════════
# SLIDE 2: Agenda / Table of Contents
# ══════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, BG_WHITE)

# Header area
add_rect(slide2, Inches(0), Inches(0), SLIDE_W, Inches(1.4), LIGHT_GRAY)
add_decorative_bar(slide2, Inches(0), Inches(0), SLIDE_W, Pt(4), PRIMARY)
add_textbox(slide2, Inches(0.8), Inches(0.35), Inches(6), Inches(0.8),
            "Agenda", font_size=36, color=TEXT_DARK, bold=True)
add_textbox(slide2, Inches(11), Inches(0.45), Inches(2), Inches(0.5),
            "QUANTROPI", font_size=11, color=PRIMARY, bold=True, alignment=PP_ALIGN.RIGHT)

# Agenda items as numbered cards
agenda_items = [
    ("01", "Topic One", "Brief description of the first agenda item"),
    ("02", "Topic Two", "Brief description of the second agenda item"),
    ("03", "Topic Three", "Brief description of the third agenda item"),
    ("04", "Topic Four", "Brief description of the fourth agenda item"),
    ("05", "Topic Five", "Brief description of the fifth agenda item"),
]

start_y = Inches(2.0)
card_h = Inches(0.85)
gap = Inches(0.15)

for i, (num, title, desc) in enumerate(agenda_items):
    y = start_y + i * (card_h + gap)
    # Number circle
    add_circle(slide2, Inches(0.9), y + Inches(0.1), Inches(0.55), PRIMARY)
    add_textbox(slide2, Inches(0.9), y + Inches(0.15), Inches(0.55), Inches(0.45),
                num, font_size=16, color=BG_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title + description
    add_textbox(slide2, Inches(1.7), y + Inches(0.05), Inches(6), Inches(0.35),
                title, font_size=18, color=TEXT_DARK, bold=True)
    add_textbox(slide2, Inches(1.7), y + Inches(0.42), Inches(8), Inches(0.35),
                desc, font_size=13, color=MID_GRAY)
    # Separator line
    if i < len(agenda_items) - 1:
        add_rect(slide2, Inches(1.7), y + card_h - Pt(1), Inches(10), Pt(1), RGBColor(0xE5, 0xE5, 0xEA))

add_footer_bar(slide2)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: Section Divider
# ══════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, TEXT_DARK)

# Decorative elements
add_circle(slide3, Inches(-1), Inches(-1), Inches(4), RGBColor(0x1A, 0x22, 0x55))
add_circle(slide3, Inches(10), Inches(4.5), Inches(5), RGBColor(0x1A, 0x22, 0x55))

# Section number
add_textbox(slide3, Inches(1.2), Inches(2.0), Inches(3), Inches(0.6),
            "01", font_size=60, color=PRIMARY, bold=True)

# Accent bar
add_decorative_bar(slide3, Inches(1.2), Inches(3.3), Inches(2), Pt(4), PRIMARY)

# Section title
add_textbox(slide3, Inches(1.2), Inches(3.6), Inches(10), Inches(1.2),
            "Section Title", font_size=40, color=BG_WHITE, bold=True)

# Section subtitle
add_textbox(slide3, Inches(1.2), Inches(4.9), Inches(8), Inches(0.6),
            "Brief description of what this section covers", font_size=18,
            color=RGBColor(0x8B, 0x90, 0xB0))

# Bottom gradient bar
add_gradient_bar(slide3, Inches(0), Inches(7.3), SLIDE_W, Pt(4))


# ══════════════════════════════════════════════════════════════
# SLIDE 4: Content Slide (Title + Body)
# ══════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, BG_WHITE)

# Top bar
add_rect(slide4, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PRIMARY)

# Small branding label
add_textbox(slide4, Inches(11), Inches(0.3), Inches(2), Inches(0.4),
            "QUANTROPI", font_size=11, color=PRIMARY, bold=True, alignment=PP_ALIGN.RIGHT)

# Slide title
add_textbox(slide4, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Content Slide Title", font_size=32, color=TEXT_DARK, bold=True)

# Accent underline
add_decorative_bar(slide4, Inches(0.8), Inches(1.2), Inches(1.5), Pt(3), ACCENT)

# Body text placeholder
add_multiline_textbox(slide4, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5),
    [
        {"text": "Key point or paragraph text goes here. Use this slide for detailed content,",
         "font_size": 18, "color": TEXT_DARK, "spacing_after": 12},
        {"text": "explanations, or narrative storytelling. Keep text concise and impactful.",
         "font_size": 18, "color": TEXT_DARK, "spacing_after": 24},
        {"text": "•  First supporting detail or bullet point", "font_size": 16,
         "color": MID_GRAY, "spacing_after": 8},
        {"text": "•  Second supporting detail or bullet point", "font_size": 16,
         "color": MID_GRAY, "spacing_after": 8},
        {"text": "•  Third supporting detail or bullet point", "font_size": 16,
         "color": MID_GRAY, "spacing_after": 8},
        {"text": "•  Fourth supporting detail or bullet point", "font_size": 16,
         "color": MID_GRAY},
    ])

add_footer_bar(slide4)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: Two-Column Layout
# ══════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, BG_WHITE)

add_rect(slide5, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PRIMARY)
add_textbox(slide5, Inches(11), Inches(0.3), Inches(2), Inches(0.4),
            "QUANTROPI", font_size=11, color=PRIMARY, bold=True, alignment=PP_ALIGN.RIGHT)

add_textbox(slide5, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Two-Column Layout", font_size=32, color=TEXT_DARK, bold=True)
add_decorative_bar(slide5, Inches(0.8), Inches(1.2), Inches(1.5), Pt(3), ACCENT)

# Left column card
add_rounded_rect(slide5, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.8), LIGHT_GRAY)
add_textbox(slide5, Inches(1.2), Inches(2.0), Inches(4.8), Inches(0.5),
            "Left Column", font_size=22, color=TEXT_DARK, bold=True)
add_decorative_bar(slide5, Inches(1.2), Inches(2.55), Inches(1), Pt(3), PRIMARY)
add_multiline_textbox(slide5, Inches(1.2), Inches(2.9), Inches(4.8), Inches(3.2),
    [
        {"text": "•  Point one for the left column", "font_size": 15, "color": MID_GRAY, "spacing_after": 8},
        {"text": "•  Point two for the left column", "font_size": 15, "color": MID_GRAY, "spacing_after": 8},
        {"text": "•  Point three for the left column", "font_size": 15, "color": MID_GRAY, "spacing_after": 8},
        {"text": "•  Add your content here", "font_size": 15, "color": MID_GRAY},
    ])

# Right column card
add_rounded_rect(slide5, Inches(6.8), Inches(1.7), Inches(5.6), Inches(4.8), LIGHT_GRAY)
add_textbox(slide5, Inches(7.2), Inches(2.0), Inches(4.8), Inches(0.5),
            "Right Column", font_size=22, color=TEXT_DARK, bold=True)
add_decorative_bar(slide5, Inches(7.2), Inches(2.55), Inches(1), Pt(3), ACCENT)
add_multiline_textbox(slide5, Inches(7.2), Inches(2.9), Inches(4.8), Inches(3.2),
    [
        {"text": "•  Point one for the right column", "font_size": 15, "color": MID_GRAY, "spacing_after": 8},
        {"text": "•  Point two for the right column", "font_size": 15, "color": MID_GRAY, "spacing_after": 8},
        {"text": "•  Point three for the right column", "font_size": 15, "color": MID_GRAY, "spacing_after": 8},
        {"text": "•  Add your content here", "font_size": 15, "color": MID_GRAY},
    ])

add_footer_bar(slide5)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: Three-Card Feature Slide
# ══════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, BG_WHITE)

add_rect(slide6, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PRIMARY)
add_textbox(slide6, Inches(11), Inches(0.3), Inches(2), Inches(0.4),
            "QUANTROPI", font_size=11, color=PRIMARY, bold=True, alignment=PP_ALIGN.RIGHT)

add_textbox(slide6, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Key Features / Highlights", font_size=32, color=TEXT_DARK, bold=True)
add_decorative_bar(slide6, Inches(0.8), Inches(1.2), Inches(1.5), Pt(3), ACCENT)

# Three cards
card_width = Inches(3.5)
card_height = Inches(4.5)
card_gap = Inches(0.5)
start_x = Inches(0.9)

colors_top = [PRIMARY, ACCENT, RGBColor(0x9C, 0xC4, 0xE4)]
titles = ["Feature One", "Feature Two", "Feature Three"]
descs = [
    "Description of this feature and why it matters to your audience.",
    "Description of this feature and why it matters to your audience.",
    "Description of this feature and why it matters to your audience.",
]
icons = ["🔒", "⚡", "🌐"]

for i in range(3):
    x = start_x + i * (card_width + card_gap)
    y = Inches(1.7)

    # Card background
    card = add_rounded_rect(slide6, x, y, card_width, card_height, LIGHT_GRAY)

    # Color top strip
    add_rect(slide6, x, y, card_width, Inches(0.15), colors_top[i])

    # Icon
    add_textbox(slide6, x + Inches(0.4), y + Inches(0.5), Inches(1), Inches(0.7),
                icons[i], font_size=36, color=TEXT_DARK, bold=False)

    # Card title
    add_textbox(slide6, x + Inches(0.4), y + Inches(1.3), Inches(2.7), Inches(0.5),
                titles[i], font_size=20, color=TEXT_DARK, bold=True)

    # Card description
    add_textbox(slide6, x + Inches(0.4), y + Inches(1.9), Inches(2.7), Inches(2),
                descs[i], font_size=14, color=MID_GRAY)

add_footer_bar(slide6)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: Data / Statistics Slide
# ══════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, BG_WHITE)

add_rect(slide7, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PRIMARY)
add_textbox(slide7, Inches(11), Inches(0.3), Inches(2), Inches(0.4),
            "QUANTROPI", font_size=11, color=PRIMARY, bold=True, alignment=PP_ALIGN.RIGHT)

add_textbox(slide7, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Key Metrics & Data", font_size=32, color=TEXT_DARK, bold=True)
add_decorative_bar(slide7, Inches(0.8), Inches(1.2), Inches(1.5), Pt(3), ACCENT)

# Four stat boxes
stats = [
    ("99.9%", "Uptime SLA"),
    ("256-bit", "Encryption"),
    ("10M+", "Keys Generated"),
    ("<1ms", "Latency"),
]

box_w = Inches(2.6)
box_h = Inches(2.8)
start_x = Inches(0.8)
box_gap = Inches(0.3)

for i, (value, label) in enumerate(stats):
    x = start_x + i * (box_w + box_gap)
    y = Inches(2.0)

    card = add_rounded_rect(slide7, x, y, box_w, box_h, LIGHT_GRAY)

    # Accent top
    add_decorative_bar(slide7, x + Inches(0.8), y + Inches(0.5), Inches(1), Pt(3),
                       PRIMARY if i % 2 == 0 else ACCENT)

    # Big number
    add_textbox(slide7, x, y + Inches(0.8), box_w, Inches(0.8),
                value, font_size=36, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    # Label
    add_textbox(slide7, x, y + Inches(1.7), box_w, Inches(0.5),
                label, font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_footer_bar(slide7)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: Quote / Highlight Slide
# ══════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, LIGHT_GRAY)

# Left accent bar
add_rect(slide8, Inches(0), Inches(0), Inches(0.15), SLIDE_H, PRIMARY)

# Quote mark
add_textbox(slide8, Inches(1.5), Inches(1.5), Inches(2), Inches(1.5),
            "\u201C", font_size=120, color=PRIMARY_LIGHT, bold=True)

# Quote text
add_textbox(slide8, Inches(1.8), Inches(2.5), Inches(9.5), Inches(2.5),
            "Insert an impactful quote, customer testimonial, or key message that reinforces your presentation narrative.",
            font_size=26, color=TEXT_DARK, bold=False)

add_decorative_bar(slide8, Inches(1.8), Inches(5.2), Inches(1.5), Pt(3), ACCENT)

add_textbox(slide8, Inches(1.8), Inches(5.5), Inches(6), Inches(0.5),
            "— Attribution or Source", font_size=16, color=MID_GRAY, bold=False)

# Branding
add_textbox(slide8, Inches(11), Inches(0.3), Inches(2), Inches(0.4),
            "QUANTROPI", font_size=11, color=PRIMARY, bold=True, alignment=PP_ALIGN.RIGHT)

add_footer_bar(slide8)


# ══════════════════════════════════════════════════════════════
# SLIDE 9: Comparison / Before-After
# ══════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, BG_WHITE)

add_rect(slide9, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PRIMARY)
add_textbox(slide9, Inches(11), Inches(0.3), Inches(2), Inches(0.4),
            "QUANTROPI", font_size=11, color=PRIMARY, bold=True, alignment=PP_ALIGN.RIGHT)

add_textbox(slide9, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Comparison", font_size=32, color=TEXT_DARK, bold=True)
add_decorative_bar(slide9, Inches(0.8), Inches(1.2), Inches(1.5), Pt(3), ACCENT)

# Left side - "Without"
add_rounded_rect(slide9, Inches(0.8), Inches(1.7), Inches(5.6), Inches(4.8), RGBColor(0xFC, 0xF0, 0xF0))
add_rect(slide9, Inches(0.8), Inches(1.7), Inches(5.6), Inches(0.6), RGBColor(0xE0, 0xC0, 0xC0))
add_textbox(slide9, Inches(1.2), Inches(1.75), Inches(4.8), Inches(0.5),
            "Without Quantropi", font_size=18, color=RGBColor(0x8B, 0x40, 0x40), bold=True)

add_multiline_textbox(slide9, Inches(1.2), Inches(2.6), Inches(4.8), Inches(3.5),
    [
        {"text": "✕  Vulnerability to quantum attacks", "font_size": 15, "color": MID_GRAY, "spacing_after": 12},
        {"text": "✕  Complex key management", "font_size": 15, "color": MID_GRAY, "spacing_after": 12},
        {"text": "✕  High latency overhead", "font_size": 15, "color": MID_GRAY, "spacing_after": 12},
        {"text": "✕  Limited scalability", "font_size": 15, "color": MID_GRAY},
    ])

# Right side - "With"
add_rounded_rect(slide9, Inches(6.8), Inches(1.7), Inches(5.6), Inches(4.8), RGBColor(0xEF, 0xF5, 0xEF))
add_rect(slide9, Inches(6.8), Inches(1.7), Inches(5.6), Inches(0.6), RGBColor(0xC0, 0xE0, 0xC0))
add_textbox(slide9, Inches(7.2), Inches(1.75), Inches(4.8), Inches(0.5),
            "With Quantropi", font_size=18, color=RGBColor(0x30, 0x7B, 0x30), bold=True)

add_multiline_textbox(slide9, Inches(7.2), Inches(2.6), Inches(4.8), Inches(3.5),
    [
        {"text": "✓  Quantum-safe encryption", "font_size": 15, "color": MID_GRAY, "spacing_after": 12},
        {"text": "✓  Automated key lifecycle", "font_size": 15, "color": MID_GRAY, "spacing_after": 12},
        {"text": "✓  Sub-millisecond performance", "font_size": 15, "color": MID_GRAY, "spacing_after": 12},
        {"text": "✓  Enterprise scale ready", "font_size": 15, "color": MID_GRAY},
    ])

add_footer_bar(slide9)


# ══════════════════════════════════════════════════════════════
# SLIDE 10: Timeline / Roadmap
# ══════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, BG_WHITE)

add_rect(slide10, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PRIMARY)
add_textbox(slide10, Inches(11), Inches(0.3), Inches(2), Inches(0.4),
            "QUANTROPI", font_size=11, color=PRIMARY, bold=True, alignment=PP_ALIGN.RIGHT)

add_textbox(slide10, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Roadmap / Timeline", font_size=32, color=TEXT_DARK, bold=True)
add_decorative_bar(slide10, Inches(0.8), Inches(1.2), Inches(1.5), Pt(3), ACCENT)

# Horizontal timeline line
add_gradient_bar(slide10, Inches(1), Inches(3.6), Inches(11.3), Pt(4))

# Timeline nodes
milestones = [
    ("Q1", "Phase 1", "Discovery &\nAssessment"),
    ("Q2", "Phase 2", "Design &\nPlanning"),
    ("Q3", "Phase 3", "Implementation\n& Testing"),
    ("Q4", "Phase 4", "Deployment &\nOptimization"),
]

node_colors = [PRIMARY, ACCENT, PRIMARY, ACCENT]

for i, (quarter, phase, desc) in enumerate(milestones):
    x_center = Inches(1.8) + i * Inches(3.0)

    # Node circle
    add_circle(slide10, x_center - Inches(0.25), Inches(3.35), Inches(0.55), node_colors[i])
    add_textbox(slide10, x_center - Inches(0.25), Inches(3.4), Inches(0.55), Inches(0.45),
                quarter, font_size=11, color=BG_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Phase label (above)
    add_textbox(slide10, x_center - Inches(1), Inches(2.4), Inches(2), Inches(0.5),
                phase, font_size=16, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)

    # Description (below)
    add_textbox(slide10, x_center - Inches(1), Inches(4.2), Inches(2), Inches(1),
                desc, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_footer_bar(slide10)


# ══════════════════════════════════════════════════════════════
# SLIDE 11: Image + Text Layout
# ══════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide11, BG_WHITE)

add_rect(slide11, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PRIMARY)
add_textbox(slide11, Inches(11), Inches(0.3), Inches(2), Inches(0.4),
            "QUANTROPI", font_size=11, color=PRIMARY, bold=True, alignment=PP_ALIGN.RIGHT)

add_textbox(slide11, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
            "Image + Text Layout", font_size=32, color=TEXT_DARK, bold=True)
add_decorative_bar(slide11, Inches(0.8), Inches(1.2), Inches(1.5), Pt(3), ACCENT)

# Left: image placeholder
img_placeholder = add_rounded_rect(slide11, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), LIGHT_GRAY)
add_textbox(slide11, Inches(1.5), Inches(3.5), Inches(4), Inches(0.8),
            "[Insert Image Here]", font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Right: text content
add_textbox(slide11, Inches(7), Inches(1.8), Inches(5.5), Inches(0.5),
            "Visual Context", font_size=24, color=TEXT_DARK, bold=True)
add_decorative_bar(slide11, Inches(7), Inches(2.4), Inches(1), Pt(3), PRIMARY)

add_multiline_textbox(slide11, Inches(7), Inches(2.8), Inches(5.5), Inches(3.5),
    [
        {"text": "Use this layout when combining visual assets with descriptive content. "
                 "The image area on the left supports photos, diagrams, screenshots, or charts.",
         "font_size": 16, "color": MID_GRAY, "spacing_after": 16},
        {"text": "•  Supporting detail one", "font_size": 15, "color": MID_GRAY, "spacing_after": 8},
        {"text": "•  Supporting detail two", "font_size": 15, "color": MID_GRAY, "spacing_after": 8},
        {"text": "•  Supporting detail three", "font_size": 15, "color": MID_GRAY},
    ])

add_footer_bar(slide11)


# ══════════════════════════════════════════════════════════════
# SLIDE 12: Call to Action / Contact
# ══════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide12, TEXT_DARK)

# Decorative circles
add_circle(slide12, Inches(8.5), Inches(-1.5), Inches(6), RGBColor(0x1A, 0x22, 0x55))
add_circle(slide12, Inches(-2), Inches(4), Inches(5), RGBColor(0x1A, 0x22, 0x55))

# Top accent
add_gradient_bar(slide12, Inches(0), Inches(0), SLIDE_W, Pt(4))

# Thank you text
add_textbox(slide12, Inches(1.5), Inches(1.5), Inches(10), Inches(0.6),
            "QUANTROPI", font_size=14, color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide12, Inches(1.5), Inches(2.3), Inches(10), Inches(1.2),
            "Thank You", font_size=48, color=BG_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_decorative_bar(slide12, Inches(5.9), Inches(3.6), Inches(1.5), Pt(4), PRIMARY)

add_textbox(slide12, Inches(2), Inches(3.9), Inches(9), Inches(0.8),
            "Questions? Let's continue the conversation.", font_size=20,
            color=RGBColor(0x8B, 0x90, 0xB0), alignment=PP_ALIGN.CENTER)

# Contact info card
add_rounded_rect(slide12, Inches(3.8), Inches(4.8), Inches(5.7), Inches(1.8),
                 RGBColor(0x1A, 0x22, 0x55))

add_multiline_textbox(slide12, Inches(4.2), Inches(4.95), Inches(4.9), Inches(1.6),
    [
        {"text": "Randy  |  Quantropi Inc.", "font_size": 16, "bold": True,
         "color": BG_WHITE, "spacing_after": 6},
        {"text": "📧  email@quantropi.com", "font_size": 13, "color": RGBColor(0x8B, 0x90, 0xB0),
         "spacing_after": 4},
        {"text": "🌐  www.quantropi.com", "font_size": 13, "color": RGBColor(0x8B, 0x90, 0xB0),
         "spacing_after": 4},
        {"text": "📍  Ottawa, Canada", "font_size": 13, "color": RGBColor(0x8B, 0x90, 0xB0)},
    ], alignment=PP_ALIGN.CENTER)

# Bottom gradient bar
add_gradient_bar(slide12, Inches(0), Inches(7.3), SLIDE_W, Pt(4))


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Randy_quantropy_template.pptx")
prs.save(output_path)
print(f"✅ Template saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Format: 16:9 Widescreen")
print(f"   Brand: Quantropi (Light Theme)")
